from __future__ import annotations

from dataclasses import dataclass
import csv
from pathlib import Path
from typing import Iterable

try:
    from pypdf import PdfReader
except Exception:  # pragma: no cover - optional dependency
    PdfReader = None  # type: ignore[assignment]

try:
    from docx import Document as DocxDocument
except Exception:  # pragma: no cover - optional dependency
    DocxDocument = None  # type: ignore[assignment]

try:
    from pptx import Presentation
except Exception:  # pragma: no cover - optional dependency
    Presentation = None  # type: ignore[assignment]


@dataclass(frozen=True)
class ExtractedDocument:
    file_id: str
    name: str
    mime_type: str | None
    local_path: str
    source_ts_utc: str | None
    extracted_text: str
    page_count: int | None
    char_count: int
    warnings: tuple[str, ...] = ()


def _clip_text(value: str, max_chars: int) -> str:
    if len(value) <= max_chars:
        return value
    return value[: max_chars - 1].rstrip() + "\n"


def _safe_read_text(path: Path) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except Exception:
            continue
    data = path.read_bytes()
    return data.decode("utf-8", errors="replace")


def _extract_pdf(path: Path) -> tuple[str, int, list[str]]:
    warnings: list[str] = []
    if PdfReader is None:
        return ("", 0, ["pdf_extractor_unavailable:pypdf_not_installed"])
    try:
        reader = PdfReader(str(path))
    except Exception as exc:
        return ("", 0, [f"pdf_parse_failed:{type(exc).__name__}"])

    chunks: list[str] = []
    pages = len(reader.pages)
    for idx, page in enumerate(reader.pages):
        try:
            text = page.extract_text() or ""
        except Exception as exc:
            warnings.append(f"pdf_page_extract_failed:{idx}:{type(exc).__name__}")
            text = ""
        if text.strip():
            chunks.append(text.strip())
    return ("\n\n".join(chunks), pages, warnings)


def _extract_docx(path: Path) -> tuple[str, int, list[str]]:
    if DocxDocument is None:
        return ("", 0, ["docx_extractor_unavailable:python-docx_not_installed"])
    try:
        doc = DocxDocument(str(path))
    except Exception as exc:
        return ("", 0, [f"docx_parse_failed:{type(exc).__name__}"])
    lines = [paragraph.text.strip() for paragraph in doc.paragraphs if paragraph.text and paragraph.text.strip()]
    return ("\n".join(lines), 1, [])


def _extract_pptx(path: Path) -> tuple[str, int, list[str]]:
    if Presentation is None:
        return ("", 0, ["pptx_extractor_unavailable:python-pptx_not_installed"])
    try:
        deck = Presentation(str(path))
    except Exception as exc:
        return ("", 0, [f"pptx_parse_failed:{type(exc).__name__}"])

    lines: list[str] = []
    for slide in deck.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if not isinstance(text, str):
                continue
            cleaned = text.strip()
            if cleaned:
                lines.append(cleaned)
    return ("\n".join(lines), len(deck.slides), [])


def _extract_csv(path: Path, *, row_limit: int = 200) -> str:
    lines: list[str] = []
    with path.open("r", encoding="utf-8", errors="replace", newline="") as handle:
        reader = csv.reader(handle)
        for idx, row in enumerate(reader):
            if idx >= row_limit:
                lines.append("[csv truncated]")
                break
            lines.append(", ".join(col.strip() for col in row))
    return "\n".join(lines)


def _extension(name: str) -> str:
    return Path(name).suffix.lower()


def _is_text_extension(ext: str) -> bool:
    return ext in {".txt", ".md", ".markdown", ".json", ".yaml", ".yml", ".log"}


def _extract_text(path: Path, *, name: str, mime_type: str | None) -> tuple[str, int | None, list[str]]:
    ext = _extension(name)
    if ext == ".pdf" or (mime_type or "").lower() == "application/pdf":
        text, pages, warnings = _extract_pdf(path)
        return (text, pages, warnings)
    if ext == ".docx" or (mime_type or "").lower() == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        text, pages, warnings = _extract_docx(path)
        return (text, pages, warnings)
    if ext == ".pptx" or (mime_type or "").lower() == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        text, pages, warnings = _extract_pptx(path)
        return (text, pages, warnings)
    if ext == ".csv" or (mime_type or "").lower() == "text/csv":
        return (_extract_csv(path), None, [])
    if _is_text_extension(ext):
        return (_safe_read_text(path), None, [])
    return ("", None, [f"unsupported_file_type:{ext or 'none'}"])


def extract_document(
    *,
    file_id: str,
    name: str,
    mime_type: str | None,
    local_path: str,
    source_ts_utc: str | None,
    max_chars: int = 120_000,
) -> ExtractedDocument:
    path = Path(local_path)
    if not path.exists() or (not path.is_file()):
        return ExtractedDocument(
            file_id=file_id,
            name=name,
            mime_type=mime_type,
            local_path=local_path,
            source_ts_utc=source_ts_utc,
            extracted_text="",
            page_count=None,
            char_count=0,
            warnings=("missing_local_path",),
        )

    text, pages, warnings = _extract_text(path, name=name, mime_type=mime_type)
    clipped = _clip_text(text, max_chars=max_chars) if text else ""
    return ExtractedDocument(
        file_id=file_id,
        name=name,
        mime_type=mime_type,
        local_path=local_path,
        source_ts_utc=source_ts_utc,
        extracted_text=clipped,
        page_count=pages,
        char_count=len(clipped),
        warnings=tuple(warnings),
    )


def extract_documents(
    rows: Iterable[dict[str, str | int | None]],
    *,
    max_chars_per_doc: int = 120_000,
) -> list[ExtractedDocument]:
    out: list[ExtractedDocument] = []
    for row in rows:
        out.append(
            extract_document(
                file_id=str(row.get("slack_file_id") or row.get("file_id") or ""),
                name=str(row.get("original_name") or row.get("name") or "unknown"),
                mime_type=(str(row.get("mimetype")) if row.get("mimetype") else None),
                local_path=str(row.get("local_path") or ""),
                source_ts_utc=(str(row.get("ingested_at_utc")) if row.get("ingested_at_utc") else None),
                max_chars=max_chars_per_doc,
            )
        )
    return out
